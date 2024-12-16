__all__ = [
    'DirectoryScanner',
    'ImgTesseractHandler',
    'HandlerBaseClass',
    'BatchingStrategyBaseClass',
    'ReporterBaseClass'
]


# imports: standard library
import hashlib
import re
import textwrap
import warnings
from abc import (
    ABC,
    ABCMeta,
    abstractmethod
)
from collections import (
    Counter,
    defaultdict
)
from collections.abc import MutableMapping
from dataclasses import (
    dataclass,
    field,
    KW_ONLY
)
from pathlib import Path
from string import (
    printable,
    Template
)
from typing import (
    Any,
    ClassVar,
    Type
)

# imports: 3rd party libraries
import black
from more_itertools import constrained_batches
from tabulate import tabulate

# imports: user modules
from sberpunk import EasyPath


# settings
def formatwarning(message, category, filename, lineno, line=None):
    return f'\n{filename}:{lineno}: {category.__name__}\n{message}\n'

warnings.formatwarning = formatwarning
warnings.filterwarnings('ignore', category=UserWarning, module='zipfile')


ENCODINGS_TO_TRY = ['utf-8', 'windows-1251']
SENTINEL_TEMPLATE = Template('[__${content}__]')
BAD_FILE_SENTINEL = SENTINEL_TEMPLATE.substitute(
    content='BAD FILE SENTINEL'
)
TH_TAG_PATTERN = re.compile(r'<th>(\w+)</th>')

HTML_TEMPLATE = Template("""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>$report_title</title>
<style>
body {
    font-family: "Segoe UI";
    background-color: #2e3440;
    color: #c9cfda;
    overflow-x: hidden;
}
a:link {
    color: #a3be8c;
}
a:visited {
    color: #b48ead;
}
table {
    border-collapse: collapse;
    width: 100%;
}
th.keywords {
    width: 15%
}
</style>
</head>
<body>
<h1>$report_title</h1>
<hr>
<h2>Table of Contents</h2>
<ol type="1">
$table_of_contents
</ol>
<hr>
<h2>Search Results</h2>
<ol type="1">
$sections
</ol>
</body>
</html>""")

TOC_ITEM_TEMPLATE = Template(
    '<li id="toc$toc_item_id"><a href="#$toc_item_id">$path</a></li>'
)

SECTION_TEMPLATE = Template(
    '<li id="$toc_item_id"><a href="file:///$path">$path</a></li>'
    '<p>$table</p>'
)


class BatchingStrategyBaseClass(ABC):
    """
    Base class of a text batching strategy.
    """

    batch_id_width: int = 5  # number of digits in batch id

    @abstractmethod
    def split(self, text: str) -> list[str]:
        """Split given text into tokens. """
        raise NotImplementedError

    @abstractmethod
    def get_len(self, token: str) -> int:
        """Return the size of a single token. """
        raise NotImplementedError

    @property
    @abstractmethod
    def batch_max_size(self) -> int:
        """Max size of a batch in the units of the token size. """
        raise NotImplementedError


class LinewiseBatchingStrategy(BatchingStrategyBaseClass):
    """
    Linewise text batching strategy.
    """

    batch_max_size = 1000
    split_pattern = re.compile(r'\n+')
    sub_pattern = re.compile(r'(?a)\s+')

    def split(self, text: str) -> list[str]:
        return self.split_pattern.split(text)

    def get_len(self, token: str) -> int:
        return len(self.sub_pattern.sub('', token))


class MSHelper:
    """
    Helper methods to handle some Microsoft file formats.
    """

    @staticmethod
    def format_worksheet_title(ws, title_attr_name: str) -> str:
        """Return XLS/XLSX worksheet title as a sentinel. """
        title = f'Worksheet {getattr(ws, title_attr_name)}'
        return SENTINEL_TEMPLATE.substitute(content=title)

    @staticmethod
    def tabulate(table: list[list[Any]]) -> str:
        """Call `tabulate.tabulate` with a list of lists `table`. """
        return tabulate(table, tablefmt='plain')

    @staticmethod
    def handle_badzipfile_error(extract_text_method):
        """
        Decorator for a text extracting method
        aimed to handle `zipfile.BadZipFile` error.
        """
        def wrapper(path):
            from zipfile import BadZipFile
            try:
                return extract_text_method(path)
            except BadZipFile as error:
                warnings.warn(f'{path}:\n{error}', UserWarning)
                return BAD_FILE_SENTINEL
        return wrapper


class HandlerBaseClass(ABC):
    """
    Base class of a file handler.
    """

    batching_strategy = None

    @staticmethod
    @abstractmethod
    def extract_text(path):
        """Extract text from a file under the given path. """
        for encoding in ENCODINGS_TO_TRY:
            try:
                with path.open(encoding=encoding) as file:
                    text = file.read()
                return text
            except UnicodeDecodeError:
                pass
        return BAD_FILE_SENTINEL

    @staticmethod
    def write_batches(extract_text_method):
        """
        Decorator for a text extracting method.
        Batch a text following the `HandlerBaseClass.batching_strategy`.
        Write batches to a TXT file created in the source file's directory.
        """

        def wrapper(path):
            strategy = HandlerBaseClass.batching_strategy
            if strategy is None:
                msg = '`HandlerBaseClass.batching_strategy` is not defined'
                raise ValueError(msg)

            text = extract_text_method(path)
            text = ''.join(c for c in text if c.isprintable() or c == '\n')
            tokens = strategy.split(text)

            batches = constrained_batches(
                iterable=tokens,
                max_size=strategy.batch_max_size,
                get_len=strategy.get_len,
                strict=False
            )

            txt_path = path.parent / f'{path.name}.txt'

            with txt_path.open('w', encoding='utf-8') as file:
                for batch_id, batch in enumerate(batches, 1):
                    width = strategy.batch_id_width
                    batch_id = str(batch_id).zfill(width)
                    batch_id = SENTINEL_TEMPLATE.substitute(
                        content=f'{batch_id=}'
                    )
                    batch = '\n'.join(batch)
                    file.write(f'{batch_id}\n{batch}\n\n\n')

            return txt_path
        return wrapper


class AnyTextFileHandler(HandlerBaseClass):
    """
    Handler of a text file with any extension.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        return HandlerBaseClass.extract_text(path)


class ArchiveHandler(HandlerBaseClass):
    """
    Handler of an archive except ZIP.
    """

    @staticmethod
    def extract_text(path):
        from patoolib import extract_archive
        new_dir = path.parent / path.stem
        new_dir.mkdir(exist_ok=True)
        extract_archive(str(path), outdir=str(new_dir))
        return new_dir


class DocxHandler(HandlerBaseClass):
    """
    DOC/DOCX document handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    @MSHelper.handle_badzipfile_error
    def extract_text(path):
        from docx2python import docx2python
        return docx2python(path).text


class ImgTesseractHandler(HandlerBaseClass):
    """
    Handler of an image.
    """

    tesseract_cmd = None
    tesseract_lang = None

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        import pytesseract
        from PIL import Image

        cmd = ImgTesseractHandler.tesseract_cmd
        lang = ImgTesseractHandler.tesseract_lang

        if cmd is None or lang is None:
            msg = (
                "`ImgTesseractHandler` expects the following kwargs:\n"
                "{'tesseract_cmd': your/path/to/tesseract.exe, "
                "'tesseract_lang': tesseract_language_code_string}.\n"
                "See https://pypi.org/project/pytesseract/ for details."
            )
            raise ValueError(msg)

        pytesseract.pytesseract.tesseract_cmd = cmd
        image = Image.open(path)
        return pytesseract.image_to_string(image, lang=lang)


class MsgHandler(HandlerBaseClass):
    """
    MSG email message handler.
    """

    @staticmethod
    def extract_text(path):
        from extract_msg import openMsg
        message = openMsg(path)
        parent_dir = path.parent
        stem = path.stem.strip()
        message.save(customPath=parent_dir, customFilename=stem)
        return parent_dir / stem


class PdfHandler(HandlerBaseClass):
    """
    PDF document handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        import struct
        from pdfminer.high_level import extract_text
        try:
            return extract_text(path)
        except struct.error as error:
            warnings.warn(f'{path}:\n{error}', UserWarning)
            return BAD_FILE_SENTINEL


class PptxHandler(HandlerBaseClass):
    """
    PPTX presentation handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        from pptx import Presentation
        text_chunks = []
        pres = Presentation(path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text_chunks.append(shape.text)
        return '\n'.join(text_chunks)


class RtfHandler(HandlerBaseClass):
    """
    RTF document handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        from striprtf.striprtf import rtf_to_text
        rtf = HandlerBaseClass.extract_text(path)
        return rtf_to_text(rtf)


class XlsHandler(HandlerBaseClass):
    """
    XLS file handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    def extract_text(path):
        from xlrd import open_workbook
        text_chunks = []
        wb = open_workbook(path)
        for ws in wb.sheets():
            title = MSHelper.format_worksheet_title(ws, 'name')
            table = [
                [str(cell.value) for cell in row]
                for row in ws.get_rows()
            ]
            text_chunks.extend([title, MSHelper.tabulate(table)])
        return '\n'.join(text_chunks)


class XlsxHandler(HandlerBaseClass):
    """
    XLSX file handler.
    """

    @staticmethod
    @HandlerBaseClass.write_batches
    @MSHelper.handle_badzipfile_error
    def extract_text(path):
        from openpyxl import load_workbook
        text_chunks = []
        wb = load_workbook(path, data_only=True)
        for ws in wb.worksheets:
            title = MSHelper.format_worksheet_title(ws, 'title')
            table = [
                list(map(XlsxHandler.materialize_cell, row))
                for row in ws.iter_rows()
            ]
            text_chunks.extend([title, MSHelper.tabulate(table)])
        return '\n'.join(text_chunks)

    @staticmethod
    def materialize_cell(cell):
        value = cell.value
        if value is None: return ''
        return value


class ZipHandler(HandlerBaseClass):
    """
    ZIP archive handler.
    """

    # zip archive is a special case:
    # you might need to fix some cyrillic filenames encoding

    _ok_chars = set(
        printable
        + (ru := 'йцукенгшщзхъфывапролджэячсмитьбюё')
        + ru.upper()
        + '№'
        + chr(8211)
        + chr(8212)
    )

    @staticmethod
    def extract_text(path):
        # (1) Write files from the original zip to a new one
        #     fixing cyrillic filenames, if needed.
        # (2) Create self titled directory and unpack new zip.
        #     Processing enters that directory.
        # (3) Finally, replace the original zip with a new zip.

        from zipfile import ZipFile, ZIP_DEFLATED

        new_dir = path.parent / path.stem
        new_dir.mkdir(exist_ok=True)
        temp_path = path.parent / 'temp.zip'

        with (
            ZipFile(path) as source,
            ZipFile(temp_path, 'w', ZIP_DEFLATED) as temp
        ):
            for name in source.namelist():
                new_name = ZipHandler._recode_filename(name)
                temp.writestr(new_name, source.read(name))

            temp.extractall(new_dir)
        temp_path.replace(path)
        return new_dir

    @staticmethod
    def _recode_filename(name: str) -> str:
        """Fix cyrillic filename encoding. """
        if not set(name).issubset(ZipHandler._ok_chars):
            try:
                return name.encode('cp437').decode('cp866')
            except UnicodeEncodeError:
                nok_chars = set(name) - ZipHandler._ok_chars
                pattern = '|'.join(map(re.escape, nok_chars))
                ZipHandler._recode_filename(re.sub(pattern, '', name))
        return name


EXTENSION_HANDLER_MAPPING = {
    '7z': ArchiveHandler,
    'csv': AnyTextFileHandler,
    'doc': DocxHandler,
    'docx': DocxHandler,
    'jpg': ImgTesseractHandler,
    'json': AnyTextFileHandler,
    'msg': MsgHandler,
    'pdf': PdfHandler,
    'png': ImgTesseractHandler,
    'pptx': PptxHandler,
    'py': AnyTextFileHandler,
    'rtf': RtfHandler,
    'txt': AnyTextFileHandler,
    'xls': XlsHandler,
    'xlsx': XlsxHandler,
    'zip': ZipHandler
}


@dataclass
class ExtensionHandlerMapping(MutableMapping):
    """
    File extension to a handler class mapping.
    """

    def __init__(self):
        self._mapping = dict(sorted(EXTENSION_HANDLER_MAPPING.items()))

    def __setitem__(self, extension, handler_class):
        self._mapping[extension] = handler_class

    def __delitem__(self, extension):
        if extension in self._mapping:
            del self._mapping[extension]

    def __getitem__(self, extension):
        return self._mapping[extension]

    def __iter__(self):
        return iter(self._mapping)

    def __len__(self):
        return len(self._mapping)

    def __str__(self):
        ext__handler_name = dict(sorted(
            (ext, handler_class.__name__)
            for ext, handler_class in self._mapping.items()
        ))
        return black.format_str(repr(ext__handler_name), mode=black.Mode())

    __repr__ = __str__


@dataclass
class Record:
    """
    Single record of the search results:
    `keywords` and the `batch` containing these `keywords`.
    """

    keywords: list[str]
    batch: str

    def __str__(self) -> str:
        attrs = {}
        attrs['keywords'] = self.keywords
        attrs['batch'] = textwrap.wrap(self.batch, width=60)
        return black.format_str(repr(attrs), mode=black.Mode())

    __repr__ = __str__


class ReporterBaseClass(ABC):
    """
    Base class of a report writer.
    """

    report_dir = None  # directory to write report to

    @abstractmethod
    def write(self, records_dict):
        """Write report with `DirectoryScanner.records` to disk. """
        raise NotImplementedError


class HtmlReporter(ReporterBaseClass):
    """
    HTML report writer.
    """

    def write(self, records_dict):
        from datetime import datetime
        from pandas import DataFrame
        from sberpunk import DATETIME_DEFAULT_FORMAT

        if self.report_dir is None:
            msg = '`HtmlReporter.report_dir` is not defined'
            raise ValueError(msg)

        table_of_contents, sections = [], []
        enumerate_obj = enumerate(records_dict.items(), 1)

        for toc_item_id, (path, records_list) in enumerate_obj:
            records = map(lambda rec: rec.__dict__, records_list)
            df = DataFrame.from_records(records)

            mapping = {
                'toc_item_id': toc_item_id,
                'path': path,
                'table': '\n' + df.to_html(justify='left', index=False)
            }

            toc_item = TOC_ITEM_TEMPLATE.substitute(mapping)
            table_of_contents.append(toc_item)

            section = SECTION_TEMPLATE.substitute(mapping)
            sections.append(section)

        now = datetime.now().strftime(DATETIME_DEFAULT_FORMAT)
        report_title = f'DirectoryScanner {now}'

        html = HTML_TEMPLATE.substitute(
            report_title=report_title,
            table_of_contents='\n'.join(table_of_contents),
            sections='<br>\n'.join(sections)
        )

        html = TH_TAG_PATTERN.sub(r'<th class="\1">\1</th>', html)

        report = self.report_dir / f'{report_title}.html'
        with report.open('w', encoding='utf-8') as file:
            file.write(html)


@dataclass
class DirectoryScanner:
    """
    Represents a core class that orchestrates successive activities of utility
    classes: recursive walking through directories, extracting text from files
    with certain extensions, searching for user defined keywords and regular
    expressions, noting the respective pieces of text and writing final report.

    Attributes:
        map (ExtensionHandlerMapping):
            (class attribute) File extension to a handler class mapping.

        keywords (list[str]):
            Keywords and/or regular expressions to search for.

        source (str | pathlib.Path):
            Source file or directory with source file(s).
            Defaults to the current working directory.

        ignore_extensions (set[str]):
            File extensions to ignore. Defaults to an empty set.

        batching_strategy (BatchingStrategyBaseClass):
            Text batching strategy.
            Defaults to `LinewiseBatchingStrategy` instance.

        handler_kwargs (dict[Type[HandlerBaseClass], dict[str, Any]]):
            Keyword arguments for the handlers. Defaults to an empty dict.

        reporter (ReporterBaseClass):
            Report writer. Defaults to `HtmlReporter` instance.
    """

    map: ClassVar[ExtensionHandlerMapping] = ExtensionHandlerMapping()

    keywords: list[str]
    _: KW_ONLY

    source: str | Path = Path.cwd()
    report_dir: str | Path = Path.cwd()

    ignore_extensions: set[str] = field(default_factory=set)
    handler_kwargs: dict[Type[HandlerBaseClass], dict[str, Any]] = field(
        default_factory=dict
    )

    batching_strategy: BatchingStrategyBaseClass = LinewiseBatchingStrategy()
    reporter: ReporterBaseClass = HtmlReporter()

    def __post_init__(self) -> None:
        instance_attrs_to_check = [
            ('batching_strategy', BatchingStrategyBaseClass),
            ('reporter', ReporterBaseClass)
        ]

        for attr_name, parent_abc in instance_attrs_to_check:
            self._check_inheritance(attr_name, parent_abc)

        HandlerBaseClass.batching_strategy = self.batching_strategy
        active_handlers = set(DirectoryScanner.map.values())

        for handler_class, handler_kwargs in self.handler_kwargs.items():
            if handler_class not in active_handlers:
                msg = (
                    f"Handler class `{handler_class.__name__}` was not "
                    f"parameterized as it's currently not responsible "
                    f"for any file extension:\n{DirectoryScanner.map}"
                )
                warnings.warn(msg, UserWarning)
                continue

            for attr, value in handler_kwargs.items():
                setattr(handler_class, attr, value)

        self.ignore_extensions = set(map(self._get_ext, self.ignore_extensions))
        self._pattern = re.compile('|'.join(self.keywords), flags=re.IGNORECASE)
        self._digests = set()
        self.records = defaultdict(list)
        self.counter = Counter()

        self.source, self.report_dir = map(Path, (self.source, self.report_dir))
        self.source = '\\\\?\\' + str(self.source.resolve())
        self._handle(self.source)

        key = lambda item: -len(item[1])
        self.records = dict(sorted(self.records.items(), key=key))
        self._write_report()

    def _handle(self, source) -> None:
        """Method managing directories and files processing. """
        for path in EasyPath(source).glob('*'):
            if path.is_dir():
                self._handle(path)

            extension = self._get_ext(path.suffix)
            if extension in self.ignore_extensions:
                continue

            handler = DirectoryScanner.map.get(extension)
            if handler is None:
                continue

            digest = self._hexdigest(path)
            if digest in self._digests:
                continue

            self._digests.add(digest)
            path = handler.extract_text(path)

            if path.is_dir():
                self._handle(path)
            else:
                self._update_records(path)

    def _update_records(self, path) -> None:
        """
        Search for the `keywords` in batches and record the search results.
        """
        with path.open(encoding='utf-8') as file:
            for batch in file.read().split('\n\n\n'):
                match_iter = self._pattern.finditer(batch)
                keywords = {m.group(0).lower() for m in match_iter}

                if keywords:
                    record = Record(sorted(keywords), batch)
                    key = str(path).removeprefix('\\\\?\\')
                    self.records[key].append(record)
                    self.counter.update(keywords)

    def _check_inheritance(
        self, attr_name: str, parent_abc: Type[ABCMeta]
    ) -> None:
        """
        Raise `ValueError` if scanner's attribute named `attr_name` is not an
        instance of a class derived from the abstract class `parent_abc`.
        """
        if not isinstance(getattr(self, attr_name), parent_abc):
            msg = (
                f'`{attr_name}` class must derive from `{parent_abc.__name__}`'
            )
            raise ValueError(msg)

    def _get_ext(self, source: str) -> str:
        """
        Get the file extension from a path suffix or clear the given extension.
        """
        return source.lstrip('*.').lower()

    def _write_report(self) -> None:
        """Write to disk a report with the search results. """
        self.reporter.report_dir = self.report_dir
        self.reporter.write(self.records)

    @staticmethod
    def _hexdigest(path) -> str:
        """Return the SHA256 digest of a file under the given path. """
        return hashlib.sha256(open(path, 'rb').read()).hexdigest()

    @classmethod
    def update_map(cls, mapping: dict[str, Type[HandlerBaseClass]]) -> None:
        """
        Update extension to a handler mapping with dict `mapping`.

        Args:
            mapping (dict[str, Type[HandlerBaseClass]]):
                Extension to a handler dict to update`DirectoryScanner.map`
                with.
        """
        for handler_class in mapping.values():
            if not issubclass(handler_class, HandlerBaseClass):
                msg = 'Handler classes must derive from `HandlerBaseClass`'
                raise ValueError(msg)
        cls.map.update(mapping)

    @classmethod
    def restore_map_defaults(cls) -> None:
        """
        Restore the original extension to a handler mapping.
        Note that this does not reset the handler kwargs.
        """
        cls.map = ExtensionHandlerMapping()

    @property
    def unique_files_seen(self) -> int:
        """Return the number of unique files seen during processing. """
        return len(self._digests)

    def __str__(self) -> str:
        return str(self.records)

    __repr__ = __str__
